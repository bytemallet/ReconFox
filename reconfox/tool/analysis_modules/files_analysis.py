import exiftool
import requests
import json
import trio
import httpx
import time
import os
import textract
from random import randint
from llama_index import Document as LlamaDocument
from django.db import IntegrityError
from reconfox.tool.ai_assistant import llama_index_helper
from reconfox.models import Domain, Files, People, Emails, Usernames, PeopleFiles
from reconfox.tool.retriever_modules import emails as emails_utils
import reconfox.tool.reconfox_utils as reconfox_utils

import sys
import rarfile
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from PyPDF2 import PdfReader
from svglib.svglib import svg2rlg
from svgwrite import Drawing

from django.db import transaction
from llama_index import VectorStoreIndex


download_direcotry = "reconfox/tool/downloaded_files/"
rate_limiter = trio.CapacityLimiter(2)


async def download_file(url, filepath):
    max_retries = 5
    file_url = url
    for i in range(max_retries + 1):
        try:
            async with httpx.AsyncClient(timeout=httpx.Timeout(60.0)) as client:
                user_agent = reconfox_utils.getUserAgents()
                response = await client.get(file_url, headers=user_agent[randint(0, len(user_agent)-1)])
                with open(filepath, 'wb') as f:
                    f.write(response.content)
                    print(f"File {url} downloaded!")
                    return 
        except Exception as e:
            print(e)
            print(type(e))
            print(f"Download failed for {url}. Retrying... ({i}/{max_retries})")
            if i%1 == 0:
                file_url = url.replace("https://", "http://")
            else:
                file_url = url              
            time.sleep(5)

async def download_all_files(domain_id):
    queryset = Files.objects.filter(metadata__isnull=True,domain_id=domain_id)
    #async with trio.open_nursery() as nursery:
    for entry in queryset.iterator():
        url = entry.url_download
        filename = entry.filename
        filepath = download_direcotry+filename
        if not os.path.isfile(os.path.join(download_direcotry, filename)):
            #async with rate_limiter:
            #task = nursery.start_soon(download_file, url, filepath)
            await download_file(url,filepath)


def downloadSingleFile(url, filename):
    print("Trying to download: " + filename)
    retry = 0
    downloaded = False
    filepath = download_direcotry+filename
    while retry < 5 and not downloaded:
        try:
            res = requests.get(url)
            open(filepath, 'wb').write(res.content)
            downloaded = True
        except Exception as e:
            retry = retry + 1
        if not downloaded:
            print(retry)
    return downloaded


def getMetadata(domain_id):
    queryset = Files.objects.filter(metadata__isnull=True,domain_id=domain_id)
    for entry in queryset.iterator():
        url = entry.url
        filename = entry.filename
        if os.path.isfile(os.path.join(download_direcotry, filename)):
            extracted = False
            retry = 0
            while not extracted and retry < 5:
                with exiftool.ExifToolHelper() as et:
                    try :
                        filepath = download_direcotry+filename
                        metadata = et.get_metadata([filepath])[0]
                        entry.metadata = metadata
                        extracted = True
                        entry.save()
                        print("metadata extracted")
                    except Exception as e:
                        print(e)
                        print(type(e))
                        retry = retry + 1
        
        


def getEmailsFromMetadata(domain_id):
    queryset = Files.objects.filter(metadata__isnull=False,domain_id=domain_id)
    for entry in queryset.iterator():
        filename = entry.filename
        metadata = str(entry.metadata)
        emails = emails_utils.getEmailsFromText(metadata)
        for e in emails:
            b,em = emails_utils.isValidEmail(e)
            if b:
                domain = Domain.objects.get(id=domain_id).domain
                if domain in em:
                    try:
                        Emails.objects.get_or_create(email=em,source="Files",domain_id=domain_id)
                    except IntegrityError as e:
                        pass



def getEmailsFromFilesContent(domain_id):
    excluded = ["rar","zip"]
    queryset = Files.objects.filter(domain_id=domain_id)
    for entry in queryset.iterator():
        url = entry.url
        filename = entry.filename
        if os.path.isfile(os.path.join(download_direcotry, filename)):
            ext = filename.split(".")[-1:][0]
            if ext not in excluded:
                #text = textract.process(os.path.join(download_direcotry, filename))
                text = extract_text(os.path.join(download_direcotry, filename))
                emails = emails_utils.getEmailsFromText(text)
                for e in emails:
                    b,em = emails_utils.isValidEmail(e)
                    if b:
                        domain = Domain.objects.get(id=domain_id).domain
                        if domain in em:
                            print("Found another email: " + em)
                            try:
                                Emails.objects.get_or_create(email=em,source="Files",domain_id=domain_id)
                            except IntegrityError as e:
                                pass
    

def getFileRelationships(domain_id):
    docs = []
    # Create a separate document for each file, including metadata extracted from JSON
    for file in Files.objects.filter(domain_id=domain_id):
        metadata = "\n".join([f"{key}: {value}" for key, value in file.metadata.items()]) if file.metadata else "None"
        doc_text = json.dumps({
            "file": file.filename,
            "metadata": file.metadata or "None"
        })
        doc = LlamaDocument(text=doc_text)
        docs.append(doc)
    
    index = VectorStoreIndex.from_documents(docs, embed_model="text-embedding-3-large")
    query_engine = index.as_query_engine()

    # Adding People, Emails, Files, and Usernames data
    for person in People.objects.filter(domain_id=domain_id):
        associated_emails = Emails.objects.filter(domain_id=domain_id, people=person).values_list('email', flat=True)
        associated_usernames = Usernames.objects.filter(domain_id=domain_id, people=person).values_list('username', flat=True)
        # Build the document text for the person
        user_data = {"person_name":person.name, "emails":list(associated_emails), "usernames":list(associated_usernames)}
        find_user_file_relationships(user_data, domain_id, query_engine)
    
    find_software(domain_id, query_engine)

        

def find_user_file_relationships(user_data, domain_id, query_engine):
    prompt = (
        f"Analyze the following user information: {user_data}. Search if the name or other details from this information are present in any file's metadata.\n"
        "Only consider a match if the metadata includes the exact name or related details from the user data provided.\n"
        "Respond in this JSON format, listing only file names where a match is found:\n"
        "{\n"
        "  \"files\": [\"<file_1>\", \"<file_2>\", \"<file_n>\"]\n"
        "}\n"
        "Example response:\n"
        "{\n"
        "  \"files\": [\"file123.json\", \"file456.json\"]\n"
        "}\n"
    )

    response = llama_index_helper.query(prompt, query_engine)
    try:
        data = json.loads(response)["files"]
        print(data)
        
        for file_name in data:
            user_name = user_data["person_name"]
            # Logic to retrieve or create the related instances
            try:
                user_name = user_data["person_name"]
                person_instance = People.objects.get(domain_id=domain_id, name=user_name)
                file_instance = Files.objects.filter(domain_id=domain_id, filename=file_name).first()
                if file_instance and person_instance:
                    PeopleFiles.objects.get_or_create(people=person_instance, file=file_instance, domain_id=domain_id)

            except People.DoesNotExist:
                print(f"Person not found: {user_name}")
            except Files.DoesNotExist:
                print(f"File not found: {file_name}")
        

    except json.JSONDecodeError as e:
        print(f"Error decoding JSON response: {e}")
        data = []

def find_software(domain_id, query_engine):
    queryset = Files.objects.filter(domain_id=domain_id)
    software_updates = []

    for file in queryset.iterator():
        prompt = (f"Give a JSON list of all software used for {file.filename}. \n"
                  "The response should follow this JSON format: {\"software\":[\"soft 1\",\"soft 2\",\"soft N\"]}")
        response = llama_index_helper.query(prompt, query_engine)

        try:
            data = json.loads(response).get("software", [])
            if isinstance(data, list) and all(isinstance(soft, str) for soft in data):
                software_updates.append((file.id, data))  # Collect updates
            else:
                print(f"Invalid data format for file {file.filename}: {data}")

        except json.JSONDecodeError as e:
            print(f"Error decoding JSON response for file {file.filename}: {e}")

    # Perform bulk update if necessary
    for file_id, software in software_updates:
        Files.objects.filter(id=file_id).update(software_used=software)






# Currently using textract
def extract_text(file_path):
    text = ""
    file_extension = file_path.split(".")[-1:][0]
    try:
        if file_extension in ["doc", "docx"]:
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"

        elif file_extension in ["ppt", "pptx", "pps", "ppsx"]:
            ppt = Presentation(file_path)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text += shape.text + "\n"

        elif file_extension in ["xls", "xlsx"]:
            wb = load_workbook(file_path)
            for sheet in wb:
                for row in sheet.iter_rows():
                    for cell in row:
                        text += str(cell.value) + " "

        elif file_extension == "pdf":
            pdf = PdfReader(file_path)
            for page_num in range(len(pdf.pages)):
                text += pdf.pages[page_num].extract_text()

        elif file_extension == "svg":
            svg = svg2rlg(file_path)
            text = str(svg)

        elif file_extension == "indd":
            print("InDesign file format (.indd) is not supported in this script.")
        
        elif file_extension == "rdp" or file_extension == "ica":
            with open(file_path, 'r') as f:
                text = f.read()

        elif file_extension == "rar":
            with rarfile.RarFile(file_path, 'r') as rar_ref:
                for file in rar_ref.namelist():
                    if not os.path.isdir(file):
                        with rar_ref.open(file) as f:
                            text += f.read().decode('utf-8', errors='ignore')

        else:
            print("Unsupported file format.")
    except Exception as e:
        print("Error occured")

    return text
